import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

DARK_BG = RGBColor(0x16, 0x16, 0x16)
CARD_BG = RGBColor(0x1F, 0x1F, 0x1F)
BORDER_COLOR = RGBColor(0x33, 0x33, 0x33)
WHITE_TEXT = RGBColor(0xF2, 0xF2, 0xF0)
GRAY_TEXT = RGBColor(0x8A, 0x8A, 0x8A)
GOLD_ACCENT = RGBColor(0xC9, 0xA2, 0x27)
RED_ACCENT = RGBColor(0xC9, 0x52, 0x4F)     # Muted brick red
GREEN_ACCENT = RGBColor(0x5F, 0xA7, 0x77)    # Muted sage green

# Typography families
FONT_HEADING = "Georgia"
FONT_BODY = "Arial"
FONT_MONO = "Courier New"

def _set_slide_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

def _add_accent_decorations(slide):
    # Add a top border accent line
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.08)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD_ACCENT
    shape.line.fill.background()

def _add_slide_header(slide, title_text):
    _add_accent_decorations(slide)
    
    # Title Text Box
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9.0), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = FONT_HEADING
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE_TEXT
    
    # Simple underline accent line
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.15), Inches(1.5), Inches(0.03)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD_ACCENT
    shape.line.fill.background()

def create_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_accent_decorations(slide)
    
    # Massive title container
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(8.0), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_HEADING
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = GOLD_ACCENT
    
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.alignment = PP_ALIGN.CENTER
        p2.font.name = FONT_BODY
        p2.font.size = Pt(20)
        p2.font.color.rgb = GRAY_TEXT
        p2.space_before = Pt(20)
        
    # Decorative line at bottom
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(5.2), Inches(3.0), Inches(0.04)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD_ACCENT
    shape.line.fill.background()

def create_content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_slide_header(slide, title)
    
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(8.6), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for idx, bullet in enumerate(bullets):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            
        p.text = "• " + bullet
        p.font.name = FONT_BODY
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE_TEXT
        p.space_after = Pt(14)

def create_two_column_slide(prs, title, left_text, right_bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_slide_header(slide, title)
    
    # Left column card shape
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.6), Inches(4.2), Inches(5.0)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = BORDER_COLOR
    card.line.width = Pt(1)
    
    leftBox = slide.shapes.add_textbox(Inches(0.85), Inches(1.85), Inches(3.7), Inches(4.5))
    ltf = leftBox.text_frame
    ltf.word_wrap = True
    
    # Handle potentially multi-line text or paragraph
    lines = left_text.split('\n')
    for idx, line in enumerate(lines):
        if idx == 0:
            p = ltf.paragraphs[0]
        else:
            p = ltf.add_paragraph()
        p.text = line
        p.font.name = FONT_BODY
        p.font.size = Pt(15)
        p.font.color.rgb = WHITE_TEXT
        p.space_after = Pt(8)
        
    # Right column bullets
    rightBox = slide.shapes.add_textbox(Inches(5.2), Inches(1.6), Inches(4.2), Inches(5.0))
    rtf = rightBox.text_frame
    rtf.word_wrap = True
    
    for idx, bullet in enumerate(right_bullets):
        if idx == 0:
            p = rtf.paragraphs[0]
        else:
            p = rtf.add_paragraph()
        p.text = "• " + bullet
        p.font.name = FONT_BODY
        p.font.size = Pt(15)
        p.font.color.rgb = WHITE_TEXT
        p.space_after = Pt(12)

def create_kpi_slide(prs, title, metrics):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_slide_header(slide, title)
    
    num_kpis = len(metrics)
    if num_kpis == 0:
        return
        
    # Calculate spacing
    total_width = 8.8  # left margin 0.6, right margin 0.6
    gap = 0.3
    card_w = (total_width - (gap * (num_kpis - 1))) / num_kpis
    
    color_map = {
        "purple": GOLD_ACCENT,
        "indigo": GOLD_ACCENT,
        "gold": GOLD_ACCENT,
        "red": RED_ACCENT,
        "green": GREEN_ACCENT
    }
    
    for idx, metric in enumerate(metrics):
        left_pos = 0.6 + idx * (card_w + gap)
        
        # Background card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left_pos), Inches(2.0), Inches(card_w), Inches(3.8)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR
        
        # Top color accent border
        color_name = metric.get("color", "indigo").lower()
        border_color = color_map.get(color_name, GOLD_ACCENT)
        
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(left_pos + 0.1), Inches(2.1), Inches(card_w - 0.2), Inches(0.08)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = border_color
        accent.line.fill.background()
        
        # Text box for KPI content
        tb = slide.shapes.add_textbox(Inches(left_pos + 0.15), Inches(2.4), Inches(card_w - 0.3), Inches(3.2))
        tf = tb.text_frame
        tf.word_wrap = True
        
        # Large Value (Mono Font)
        p_val = tf.paragraphs[0]
        p_val.text = str(metric.get("value", "0"))
        p_val.alignment = PP_ALIGN.CENTER
        p_val.font.name = FONT_MONO
        p_val.font.size = Pt(54)
        p_val.font.bold = True
        p_val.font.color.rgb = border_color
        p_val.space_after = Pt(14)
        
        # Label
        p_lbl = tf.add_paragraph()
        p_lbl.text = metric.get("label", "Metric")
        p_lbl.alignment = PP_ALIGN.CENTER
        p_lbl.font.name = FONT_BODY
        p_lbl.font.size = Pt(15)
        p_lbl.font.color.rgb = WHITE_TEXT
