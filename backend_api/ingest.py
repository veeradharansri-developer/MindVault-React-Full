import os
import pickle
import numpy as np
import io
import csv
from utils.encoder import HashedTfidfEncoder

# Multi-format parsing libraries with safe imports
try:
    from pypdf import PdfReader
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import openpyxl
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

def chunk_text(text, max_chunk_size=1000):
    """
    Split text into paragraphs, merging them into chunks of roughly max_chunk_size characters.
    """
    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
    chunks = []
    current_chunk = []
    current_length = 0
    
    for para in paragraphs:
        para_length = len(para)
        if current_length + para_length > max_chunk_size and current_chunk:
            chunks.append("\n\n".join(current_chunk))
            current_chunk = [para]
            current_length = para_length
        else:
            current_chunk.append(para)
            current_length += para_length + 2
            
    if current_chunk:
        chunks.append("\n\n".join(current_chunk))
        
    return chunks

def extract_text_and_pages(file_bytes: bytes, filename: str) -> list[dict]:
    """
    Extracts text and associates page numbers for supported file formats.
    Returns list of dicts: [{"text": str, "page_number": int | None}]
    """
    ext = os.path.splitext(filename)[1].lower()
    
    # 1. PDF
    if ext == ".pdf":
        if not PDF_AVAILABLE:
            return [{"text": "PDF parser unavailable.", "page_number": None}]
        chunks = []
        try:
            reader = PdfReader(io.BytesIO(file_bytes))
            for page_idx, page in enumerate(reader.pages):
                page_text = page.extract_text() or ""
                page_text = page_text.strip()
                if page_text:
                    p_chunks = chunk_text(page_text)
                    for c in p_chunks:
                        chunks.append({"text": c, "page_number": page_idx + 1})
        except Exception as e:
            print(f"Error parsing PDF {filename}: {e}")
        return chunks
        
    # 2. DOCX
    elif ext == ".docx":
        if not DOCX_AVAILABLE:
            return [{"text": "DOCX parser unavailable.", "page_number": None}]
        chunks = []
        try:
            doc = docx.Document(io.BytesIO(file_bytes))
            text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
            doc_chunks = chunk_text(text)
            for c in doc_chunks:
                chunks.append({"text": c, "page_number": 1})
        except Exception as e:
            print(f"Error parsing DOCX {filename}: {e}")
        return chunks
        
    # 3. PPTX
    elif ext == ".pptx":
        if not PPTX_AVAILABLE:
            return [{"text": "PPTX parser unavailable.", "page_number": None}]
        chunks = []
        try:
            prs = Presentation(io.BytesIO(file_bytes))
            for slide_idx, slide in enumerate(prs.slides):
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                slide_text = slide_text.strip()
                if slide_text:
                    p_chunks = chunk_text(slide_text)
                    for c in p_chunks:
                        chunks.append({"text": c, "page_number": slide_idx + 1})
        except Exception as e:
            print(f"Error parsing PPTX {filename}: {e}")
        return chunks
        
    # 4. XLSX
    elif ext == ".xlsx":
        if not XLSX_AVAILABLE:
            return [{"text": "XLSX parser unavailable.", "page_number": None}]
        chunks = []
        try:
            wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
            for sheet_idx, sheet_name in enumerate(wb.sheetnames):
                sheet_text = f"Sheet: {sheet_name}\n"
                ws = wb[sheet_name]
                for row in ws.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) for cell in row if cell is not None])
                    if row_text.strip():
                        sheet_text += row_text + "\n"
                sheet_text = sheet_text.strip()
                if sheet_text:
                    p_chunks = chunk_text(sheet_text)
                    for c in p_chunks:
                        chunks.append({"text": c, "page_number": sheet_idx + 1})
        except Exception as e:
            print(f"Error parsing XLSX {filename}: {e}")
        return chunks

    # 5. CSV
    elif ext == ".csv":
        chunks = []
        try:
            content = file_bytes.decode("utf-8", errors="ignore")
            reader = csv.reader(io.StringIO(content))
            text = "\n".join([", ".join(row) for row in reader if row])
            csv_chunks = chunk_text(text)
            for c in csv_chunks:
                chunks.append({"text": c, "page_number": 1})
        except Exception as e:
            print(f"Error parsing CSV {filename}: {e}")
        return chunks
        
    # 6. TXT / Default
    else:
        chunks = []
        try:
            text = file_bytes.decode("utf-8", errors="ignore")
            txt_chunks = chunk_text(text)
            for c in txt_chunks:
                chunks.append({"text": c, "page_number": 1})
        except Exception as e:
            print(f"Error parsing TXT {filename}: {e}")
        return chunks

def run_ingestion(hr_docs_dir=None, output_dir=None, model_name='all-MiniLM-L6-v2'):
    """
    Loads all supported files in hr_docs_dir, chunks them, embeds them,
    and writes metadata.pkl and embeddings.npy to output_dir.
    """
    base_dir = os.path.dirname(os.path.abspath(__file__))
    if hr_docs_dir is None:
        hr_docs_dir = os.path.join(base_dir, "hr_docs")
    if output_dir is None:
        output_dir = os.path.join(base_dir, "data")
        
    os.makedirs(output_dir, exist_ok=True)
    
    if not os.path.exists(hr_docs_dir):
        print(f"Error: {hr_docs_dir} directory does not exist.")
        return False
        
    supported_extensions = {".txt", ".pdf", ".docx", ".pptx", ".xlsx", ".csv"}
    chunks_metadata = []
    
    for filename in os.listdir(hr_docs_dir):
        ext = os.path.splitext(filename)[1].lower()
        if ext in supported_extensions:
            filepath = os.path.join(hr_docs_dir, filename)
            try:
                with open(filepath, "rb") as f:
                    file_bytes = f.read()
                
                doc_chunks = extract_text_and_pages(file_bytes, filename)
                for idx, item in enumerate(doc_chunks):
                    chunks_metadata.append({
                        "text": item["text"],
                        "source": filename,
                        "chunk_index": idx,
                        "page_number": item["page_number"]
                    })
            except Exception as e:
                print(f"Error reading file {filename}: {e}")
                
    if not chunks_metadata:
        print("No document chunks found to ingest.")
        return False
        
    print(f"Total chunks created: {len(chunks_metadata)}")
    
    model = HashedTfidfEncoder(dimension=384)
    print("Generating embeddings...")
    chunk_texts = [chunk["text"] for chunk in chunks_metadata]
    embeddings = model.encode(chunk_texts, show_progress_bar=True, convert_to_numpy=True)
    
    embeddings_file = os.path.join(output_dir, "embeddings.npy")
    metadata_file = os.path.join(output_dir, "metadata.pkl")
    
    np.save(embeddings_file, embeddings)
    with open(metadata_file, "wb") as f:
        pickle.dump(chunks_metadata, f)
        
    print(f"Successfully saved embeddings to {embeddings_file}")
    print(f"Successfully saved metadata to {metadata_file}")
    return True

if __name__ == "__main__":
    run_ingestion()
